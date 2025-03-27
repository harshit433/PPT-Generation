from fastapi import FastAPI, Form
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Pt
from text_processor import TextProcessor

app = FastAPI()

@app.post("/create-ppt/")
async def create_ppt(
    title: str = Form(...),
    content: str = Form(...),
    output_filename: str = Form(...),
    max_chars_per_slide: int = Form(500)
):
    """
    Endpoint to create a PowerPoint file with automatic slide splitting for long content.
    
    Args:
        title: Title slide text
        content: Main content text
        output_filename: Name for the output file
        max_chars_per_slide: Maximum characters per slide (default: 500)
    Returns:
        Downloadable PPT file
    """
    # Initialize text processor
    text_processor = TextProcessor(max_chars_per_slide=max_chars_per_slide)
    
    # Split content into chunks
    content_chunks = text_processor.split_into_chunks(content)
    
    # Create a new PowerPoint presentation
    presentation = Presentation()
    
    # Add a title slide
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    # Set title slide content
    title_shape.text = title
    subtitle_shape.text = f"Generated presentation with {len(content_chunks)} content slides"
    
    # Add content slides
    for i, chunk in enumerate(content_chunks, 1):
        content_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(content_slide_layout)
        
        # Set slide title
        title_shape = slide.shapes.title
        title_shape.text = f"{title} - Part {i}"
        
        # Format and set content
        content_shape = slide.placeholders[1]
        formatted_chunk = text_processor.format_chunk_for_slide(chunk)
        content_shape.text = formatted_chunk
        
        # Adjust text size if needed
        text_frame = content_shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)  # Default size
            if len(chunk) > 250:  # Reduce font size for longer chunks
                paragraph.font.size = Pt(14)
    
    # Save the presentation
    output_file = f"{output_filename}.pptx"
    presentation.save(output_file)
    
    # Return the file for download
    return FileResponse(
        output_file,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=output_file
    )