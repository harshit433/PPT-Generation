from pptx import Presentation

ppt = Presentation('D:/Harshit/Ibind Systems/PPT generation/trial.pptx')

for slide_index, slide in enumerate(ppt.slides):
    print(f"Slide {slide_index + 1}:")
    
    # Iterate over each shape on the slide
    for shape_index, shape in enumerate(slide.shapes):
        # Check if the shape has a text frame
        if shape.has_text_frame:
            # Print the text content of the shape
            print(f"  Shape {shape_index + 1}: {shape.text_frame.text}")
        elif shape.shape_type == 6:  # GroupShape
            # If the shape is a group, iterate over its sub-shapes
            for sub_shape_index, sub_shape in enumerate(shape.shapes):
                if sub_shape.has_text_frame:
                    print(f"    Group Shape {shape_index + 1}, Sub-Shape {sub_shape_index + 1}: {sub_shape.text_frame.text}")

    print()