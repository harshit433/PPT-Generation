from pptx import Presentation
import os

def update_text_preserving_format(shape, new_text):
    if not shape.has_text_frame:
        return  # Skip shapes without a text frame

    text_frame = shape.text_frame
    paragraphs = text_frame.paragraphs

    # If there are no paragraphs or text, we can't preserve formatting
    if not paragraphs or len(paragraphs[0].runs) == 0:
        text_frame.text = new_text  # Set new text directly
        return

    # Collect formatting from the first run
    first_run = paragraphs[0].runs[0]
    font = first_run.font

    # Store formatting details
    formatting = {
        "bold": font.bold,
        "italic": font.italic,
        "underline": font.underline,
        "color": font.color.rgb if font.color.type == 1 else None,
        "name": font.name,
    }

    # Clear existing text
    text_frame.clear()

    # Add new text
    p = text_frame.add_paragraph()
    run = p.add_run()
    run.text = new_text

    # Apply formatting
    run.font.bold = formatting["bold"]
    run.font.italic = formatting["italic"]
    run.font.underline = formatting["underline"]
    if formatting["color"]:
        run.font.color.rgb = formatting["color"]
    run.font.name = formatting["name"]


def modify_ppt(ppt_content, ppt_path):
    # Load the PowerPoint presentation
    presentation = Presentation(os.path.abspath(ppt_path))

    # Iterate through ppt_content and corresponding slides
    for idx, slide in zip(ppt_content['slides'], presentation.slides):
        layout_id = idx['layout_id']

        if layout_id == 1:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.shape_id == 3:
                    # shape.text = idx["content"].get('title', "")
                    update_text_preserving_format(shape, str.upper(idx["content"].get('title', "")))
                if hasattr(shape, "text") and shape.shape_id == 4:
                    # shape.text = idx["content"].get('subtitle', "")
                    update_text_preserving_format(shape, idx["content"].get('subtitle', ""))
        elif layout_id == 2:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.shape_id == 6:
                    # shape.text = idx["content"].get('heading', "")
                    update_text_preserving_format(shape, str.upper(idx["content"].get('heading', "")))
                if hasattr(shape, "text") and shape.shape_id == 7:
                    # shape.text = idx["content"].get('text_content', "")
                    update_text_preserving_format(shape, idx["content"].get('text_content', ""))
        elif layout_id == 3:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.shape_id == 3:
                    # shape.text = idx["content"].get('heading', "")
                    update_text_preserving_format(shape, str.upper(idx["content"].get('heading', "")))
                if hasattr(shape, "text") and shape.shape_id == 10:
                    # shape.text = idx["content"].get('text_content_1', "")
                    update_text_preserving_format(shape, idx["content"].get('text_content_1', ""))
                if hasattr(shape, "text") and shape.shape_id == 11:
                    # shape.text = idx["content"].get('text_content_2', "")
                    update_text_preserving_format(shape, idx["content"].get('text_content_2', ""))
        elif layout_id == 4:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.shape_id == 3:
                    # shape.text = idx["content"].get('heading', "")
                    update_text_preserving_format(shape,str.upper( idx["content"].get('heading', "")))
                if hasattr(shape, "text") and shape.shape_id == 7:
                    # shape.text = idx["content"].get('text_content_1', "")
                    update_text_preserving_format(shape, idx["content"].get('text_content_1', ""))
                if hasattr(shape, "text") and shape.shape_id == 11:
                    # shape.text = idx["content"].get('text_content_2', "")
                    update_text_preserving_format(shape, idx["content"].get('text_content_2', ""))
                if hasattr(shape, "text") and shape.shape_id == 15:
                    # shape.text = idx["content"].get('text_content_3', "")
                    update_text_preserving_format(shape, idx["content"].get('text_content_3', ""))
        elif layout_id == 5:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.shape_id == 3:
                    # shape.text = idx["content"].get('thank_you_text', "")
                    update_text_preserving_format(shape, str.upper(idx["content"].get('thank_you_text', "")))

    presentation.save(os.path.abspath(r"final_result.pptx"))         